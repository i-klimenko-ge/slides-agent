import os

try:
    from pptx import Presentation as PptxPresentation
except Exception:  # pragma: no cover - library might be missing
    PptxPresentation = None

try:
    from PyPDF2 import PdfReader
except Exception:  # pragma: no cover - library might be missing
    PdfReader = None


class PresentationReader:
    """Unified reader for PPTX and PDF presentations."""

    def __init__(self, path: str):
        self.path = path
        self.reader = self._create_reader(path)
        self.ext = os.path.splitext(path)[1].lower()

    def _create_reader(self, path: str):
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pptx":
            if PptxPresentation is None:
                raise ImportError("python-pptx is required for PPTX files")
            return PptxPresentation(path)
        if ext == ".pdf":
            if PdfReader is None:
                raise ImportError("PyPDF2 is required for PDF files")
            return PdfReader(path)
        raise ValueError(f"Unsupported file type: {ext}")

    @property
    def slides_count(self) -> int:
        if self.ext == ".pptx":
            return len(self.reader.slides)
        if self.ext == ".pdf":
            return len(self.reader.pages)
        return 0

    def get_slide_text(self, slide_number: int) -> str:
        if self.ext == ".pptx":
            slide = self.reader.slides[slide_number - 1]
            return "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
        if self.ext == ".pdf":
            page = self.reader.pages[slide_number - 1]
            text = page.extract_text()
            return text if text is not None else ""
        return ""
