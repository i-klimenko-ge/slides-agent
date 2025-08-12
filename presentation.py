"""Presentation abstractions for PPTX and PDF files with viewer support."""

from __future__ import annotations

import os
from abc import ABC, abstractmethod

from viewer import PresentationViewer

try:
    from pptx import Presentation as PptxLib
except Exception:  # pragma: no cover - optional dependency
    PptxLib = None

# optional pdf reader
try:
    from PyPDF2 import PdfReader
except Exception:  # pragma: no cover - optional dependency
    PdfReader = None


class BasePresentation(ABC):
    def __init__(self, path: str, viewer: PresentationViewer) -> None:
        self.path = path
        self.viewer = viewer

    def open(self) -> None:
        self.viewer.open(self.path)

    def close(self) -> None:
        self.viewer.close()

    def start_show(self) -> None:
        """Start presentation in fullscreen mode."""
        self.viewer.start_show()

    def goto(self, num: int) -> None:
        self.viewer.goto_slide(num)

    def next_slide(self) -> None:
        """Move to the next slide."""
        self.viewer.next_slide()

    def previous_slide(self) -> None:
        """Move to the previous slide."""
        self.viewer.previous_slide()

    @abstractmethod
    def slides_count(self) -> int:  # pragma: no cover - interface
        ...

    @abstractmethod
    def get_slide_text(self, num: int) -> str:  # pragma: no cover - interface
        ...


class PptxPresentation(BasePresentation):
    def __init__(self, path: str, viewer: PresentationViewer) -> None:
        if PptxLib is None:
            raise RuntimeError("pptx library is not available")
        super().__init__(path, viewer)
        self.prs = PptxLib(path)

    def slides_count(self) -> int:
        return len(self.prs.slides)

    def get_slide_text(self, num: int) -> str:
        slide = self.prs.slides[num]
        return "\n".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )


class PdfPresentation(BasePresentation):
    def __init__(self, path: str, viewer: PresentationViewer) -> None:
        super().__init__(path, viewer)
        if PdfReader is not None:
            try:
                self.reader = PdfReader(path)
            except Exception:
                self.reader = None
        else:
            self.reader = None

    def slides_count(self) -> int:
        if self.reader is not None:
            return len(self.reader.pages)
        return 0

    def get_slide_text(self, num: int) -> str:
        if self.reader is None:
            return ""
        try:
            page = self.reader.pages[num]
            text = page.extract_text()
            return text or ""
        except Exception:
            return ""


def create_presentation(path: str, viewer: PresentationViewer) -> BasePresentation:
    _, ext = os.path.splitext(path.lower())
    if ext == ".pptx":
        return PptxPresentation(path, viewer)
    if ext == ".pdf":
        return PdfPresentation(path, viewer)
    raise ValueError(f"Unsupported file type: {ext}")
