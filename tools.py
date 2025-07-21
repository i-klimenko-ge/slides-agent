from langchain_core.tools import tool
from typing import Annotated
import os
import platform
from pptx import Presentation as PresentationFactory
from pptx.presentation import Presentation

from config import config

PRESENTATIONS_DIR = config.get("presentations_dir", "presentations")
OS_TYPE = config.get("os", platform.system().lower())

_current_presentation: Presentation | None = None
_current_presentation_path: str | None = None
_current_slide_index: int | None = None

from viewer import get_viewer

_viewer = get_viewer(OS_TYPE)

@tool
def list_presentations_tool() -> dict:
    """Получить список файлов презентаций в каталоге."""
    if not os.path.isdir(PRESENTATIONS_DIR):
        return {
            "status": "error",
            "message": f"Каталог {PRESENTATIONS_DIR} не найден"
        }
    files = [
        f for f in os.listdir(PRESENTATIONS_DIR)
        if os.path.isfile(os.path.join(PRESENTATIONS_DIR, f))
    ]
    return {"status": "ok", "files": files}

@tool
def open_presentation_tool(query: Annotated[str, "имя файла презентации"]) -> dict:
    """Открыть презентацию для просмотра"""
    global _current_presentation, _current_presentation_path, _current_slide_index

    path = query
    if not os.path.isabs(path):
        path = os.path.join(PRESENTATIONS_DIR, path)

    if not os.path.exists(path):
        return {"status": "error", "message": f"Файл {query} не найден"}

    try:
        prs = PresentationFactory(path)
    except Exception as e:  # pragma: no cover - basic error reporting
        return {"status": "error", "message": f"Не удалось открыть файл: {e}"}

    # try to open the presentation with a system viewer if possible
    try:
        _viewer.open(path)
    except Exception:
        pass

    _current_presentation = prs
    _current_presentation_path = path
    _current_slide_index = 0
    return {
        "status": "ok",
        "slides_count": len(prs.slides),
        "message": f"Открыта презентация {os.path.basename(path)}",
    }

@tool
def close_presentation_tool() -> dict:
    """Завершить просмотр и закрыть презентацию"""
    global _current_presentation, _current_presentation_path, _current_slide_index

    if _current_presentation is None:
        return {"status": "error", "message": "Презентация не открыта"}

    _current_presentation = None
    _current_presentation_path = None
    _current_slide_index = None

    _viewer.close()

    return {"status": "ok", "message": "Презентация закрыта"}

@tool
def open_slide(slide_number: Annotated[int, "номер слайда"]) -> dict:
    """Открыть необходимый слайд в презентации по его номеру"""
    global _current_presentation, _current_slide_index

    if _current_presentation is None:
        return {"status": "error", "message": "Презентация не открыта"}

    prs = _current_presentation

    if slide_number < 1 or slide_number > len(prs.slides):
        return {"status": "error", "message": "Некорректный номер слайда"}

    slide = prs.slides[slide_number - 1]

    try:
        _viewer.goto_slide(slide_number - 1)
    except Exception:
        pass

    _current_slide_index = slide_number - 1
    text = "\n".join(
        shape.text for shape in slide.shapes if hasattr(shape, "text")
    )

    return {"status": "ok", "slide_number": slide_number, "text": text}
